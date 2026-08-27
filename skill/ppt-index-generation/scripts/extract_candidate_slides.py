#!/usr/bin/env python3
"""Extract selected slides from one PPTX into a new candidate pack."""

from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path

from pptx import Presentation


def parse_slide_list(value: str) -> list[int]:
    slides: list[int] = []
    for part in value.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            bounds = [item.strip() for item in part.split("-", 1)]
            if len(bounds) != 2 or not all(item.isdigit() for item in bounds):
                raise ValueError(f"Invalid slide range: {part}")
            start, end = map(int, bounds)
            if start > end:
                raise ValueError(f"Slide range must be ascending: {part}")
            slides.extend(range(start, end + 1))
        elif part.isdigit():
            slides.append(int(part))
        else:
            raise ValueError(f"Invalid slide number: {part}")

    slides = list(dict.fromkeys(slides))
    if not slides or any(slide < 1 for slide in slides):
        raise ValueError("At least one positive slide number is required.")
    return slides


def first_text(slide) -> str:
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.replace("\n", " ").strip()
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return shape.text.replace("\n", " ").strip()[:120]
    return "Untitled"


def file_state(path: Path) -> tuple[int, int]:
    stat = path.stat()
    return stat.st_size, stat.st_mtime_ns


def validate_paths(source: Path, output: Path, notes: Path | None) -> None:
    if not source.exists() or not source.is_file():
        raise FileNotFoundError(source)
    if source.suffix.lower() != ".pptx":
        raise ValueError(f"Only .pptx sources are supported: {source}")
    if output.suffix.lower() != ".pptx":
        raise ValueError(f"Output must use the .pptx extension: {output}")
    if source == output:
        raise ValueError("Output path must not be the source PPTX path.")
    if notes and notes in {source, output}:
        raise ValueError("Notes path must not overwrite the source or output PPTX.")


def extract(source: Path, slides: list[int], output: Path, notes: Path | None, title: str) -> None:
    source = source.expanduser().resolve()
    output = output.expanduser().resolve()
    notes = notes.expanduser().resolve() if notes else None
    validate_paths(source, output, notes)
    source_state = file_state(source)

    prs = Presentation(str(source))
    slide_count = len(prs.slides)
    invalid = [num for num in slides if num > slide_count]
    if invalid:
        raise ValueError(f"Invalid slide numbers for {source}: {invalid}")

    slide_id_list = prs.slides._sldIdLst
    original_ids = list(slide_id_list)
    selected_ids = [original_ids[num - 1] for num in slides]
    for slide_id in original_ids:
        slide_id_list.remove(slide_id)
    for slide_id in selected_ids:
        slide_id_list.append(slide_id)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    if file_state(source) != source_state:
        raise RuntimeError(f"Source PPTX changed during extraction: {source}")

    check = Presentation(str(output))
    if len(check.slides) != len(slides):
        raise RuntimeError("Generated PPTX has an unexpected slide count.")

    if notes:
        notes.parent.mkdir(parents=True, exist_ok=True)
        lines = [
            f"# {title}",
            "",
            f"- generated_at: {datetime.now().isoformat(timespec='seconds')}",
            f"- output_pptx: `{output}`",
            f"- source_pptx: `{source}`",
            f"- source_slides: {', '.join(map(str, slides))}",
            f"- output_slide_count: {len(check.slides)}",
            "",
            "## Page Map",
        ]
        for new_num, old_num in enumerate(slides, start=1):
            lines.append(f"- New slide {new_num}: source slide {old_num} - {first_text(check.slides[new_num - 1])}")
        notes.write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", required=True, type=Path, help="Source PPTX path.")
    parser.add_argument("--slides", required=True, help="Slides to keep, e.g. 1,5,7 or 23-28.")
    parser.add_argument("--output", required=True, type=Path, help="Output PPTX path.")
    parser.add_argument("--notes", type=Path, help="Optional Markdown page-map notes path.")
    parser.add_argument("--title", default="Candidate Slide Pack", help="Notes title.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    extract(args.source, parse_slide_list(args.slides), args.output, args.notes, args.title)
    print(args.output.expanduser().resolve())


if __name__ == "__main__":
    main()
