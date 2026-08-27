#!/usr/bin/env python3
"""Build AI-readable Markdown indexes for user-provided PPTX files."""

from __future__ import annotations

import argparse
from collections import Counter
from datetime import datetime
import hashlib
import json
from pathlib import Path
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


LATIN_STOPWORDS = {
    "and",
    "are",
    "for",
    "from",
    "into",
    "that",
    "the",
    "this",
    "with",
    "your",
}

COMMON_CN_STOPWORDS = {
    "一个",
    "以及",
    "进行",
    "通过",
    "实现",
    "提供",
    "基于",
    "相关",
    "主要",
    "内容",
    "工作",
}


def clean_text(text: str) -> str:
    text = text.replace("\x0b", "\n").replace("\r", "\n")
    lines = []
    for line in text.split("\n"):
        line = re.sub(r"\s+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def safe_name(name: str) -> str:
    return re.sub(r'[\\/:*?"<>|]+', "_", name).strip() or "deck"


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def extract_text_from_shape(shape) -> list[str]:
    chunks: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = clean_text(shape.text)
        if text:
            chunks.append(text)
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            cells = [clean_text(cell.text) for cell in row.cells]
            row_text = " | ".join(cell for cell in cells if cell)
            if row_text:
                rows.append(row_text)
        if rows:
            chunks.append("\n".join(rows))
    return chunks


def shape_counts(slide) -> dict[str, int]:
    counts = {"text": 0, "table": 0, "chart": 0, "image": 0, "group": 0, "other": 0}
    for shape in iter_shapes(slide.shapes):
        try:
            if getattr(shape, "has_table", False):
                counts["table"] += 1
            elif getattr(shape, "has_chart", False):
                counts["chart"] += 1
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                counts["image"] += 1
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                counts["group"] += 1
            elif getattr(shape, "has_text_frame", False) and clean_text(shape.text):
                counts["text"] += 1
            else:
                counts["other"] += 1
        except (AttributeError, TypeError, ValueError):
            counts["other"] += 1
    return {key: value for key, value in counts.items() if value}


def get_title(slide, texts: list[str]) -> str:
    try:
        if slide.shapes.title and clean_text(slide.shapes.title.text):
            return clean_text(slide.shapes.title.text).split("\n")[0][:120]
    except (AttributeError, ValueError):
        pass
    for text in texts:
        for line in text.split("\n"):
            line = line.strip()
            if 2 <= len(line) <= 100:
                return line[:120]
    return "Untitled slide"


def extract_keywords(text: str, limit: int = 6) -> list[str]:
    """Extract generic keywords without assuming an industry vocabulary."""
    latin = [
        token.lower()
        for token in re.findall(r"[A-Za-z][A-Za-z0-9_-]{2,}", text)
        if token.lower() not in LATIN_STOPWORDS
    ]
    chinese = [
        token
        for token in re.findall(r"[\u4e00-\u9fff]{2,8}", text)
        if token not in COMMON_CN_STOPWORDS
    ]

    counts = Counter(latin)
    for token in chinese:
        counts[token] += 2 if len(token) >= 4 else 1

    keywords = [token for token, _count in counts.most_common(limit)]
    return keywords or ["review-needed"]


def content_type(counts: dict[str, int]) -> str:
    parts = []
    if counts.get("text"):
        parts.append("text")
    if counts.get("table"):
        parts.append("table")
    if counts.get("chart"):
        parts.append("chart")
    if counts.get("image"):
        parts.append("image")
    if counts.get("group"):
        parts.append("diagram/group")
    return " + ".join(parts) if parts else "unknown"


def one_line(text: str, limit: int = 320) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit] + ("..." if len(text) > limit else "")


def validate_source(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(path)
    if not path.is_file():
        raise ValueError(f"PPTX path is not a file: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Only .pptx files are supported: {path}")


def build_index(paths: list[Path], output_dir: Path) -> None:
    indexes_dir = output_dir / "indexes"
    indexes_dir.mkdir(parents=True, exist_ok=True)
    master = []

    for path in paths:
        validate_source(path)
        prs = Presentation(str(path))
        deck_id = hashlib.sha256(str(path).encode("utf-8")).hexdigest()[:10]
        index_path = indexes_dir / f"{safe_name(path.stem)}--{deck_id}.md"
        stat = path.stat()
        slides = []
        deck_texts = []

        for num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in iter_shapes(slide.shapes):
                texts.extend(extract_text_from_shape(shape))
            dedup = list(dict.fromkeys(texts))
            full_text = "\n".join(dedup)
            deck_texts.append(full_text)
            counts = shape_counts(slide)
            slides.append(
                {
                    "num": num,
                    "title": get_title(slide, dedup),
                    "text": full_text,
                    "summary": one_line(full_text)
                    if full_text
                    else "No extractable text; this may be an image-only or decorative slide.",
                    "tags": extract_keywords(f"{path.stem}\n{full_text}"),
                    "content_type": content_type(counts),
                    "shape_counts": counts,
                }
            )

        deck_tags = extract_keywords(f'{path.stem}\n{" ".join(deck_texts)}')
        lines = [
            f"# {path.stem}",
            "",
            "## Deck Metadata",
            f"- deck_id: `{deck_id}`",
            f"- source_path: `{path}`",
            f"- file_size_mb: {stat.st_size / 1024 / 1024:.1f}",
            f"- modified_time: {datetime.fromtimestamp(stat.st_mtime).isoformat(timespec='seconds')}",
            f"- slide_count: {len(slides)}",
            f"- deck_tags: {', '.join(deck_tags)}",
            "- index_version: v0.2-text-structure",
            "",
            "## Slides",
            "",
        ]
        for slide in slides:
            lines.extend(
                [
                    f"### Slide {slide['num']:02d}: {slide['title']}",
                    f"- source: `{path}#slide={slide['num']}`",
                    f"- content_type: {slide['content_type']}",
                    f"- reusable_for: {', '.join(slide['tags'])}",
                    f"- summary: {slide['summary']}",
                    f"- shape_counts: `{json.dumps(slide['shape_counts'], ensure_ascii=False)}`",
                ]
            )
            if slide["text"]:
                lines.append("- extracted_text: |")
                text_lines = slide["text"].split("\n")
                for line in text_lines[:80]:
                    lines.append(f"    {line}")
                if len(text_lines) > 80:
                    lines.append("    ...")
            else:
                lines.append("- extracted_text: none")
            lines.append("")
        index_path.write_text("\n".join(lines), encoding="utf-8")
        master.append(
            (path, index_path, deck_id, stat, len(slides), deck_tags, [slide["title"] for slide in slides[:8]])
        )

    lines = [
        "# PPT Material Master Index",
        "",
        f"- generated_at: {datetime.now().isoformat(timespec='seconds')}",
        "- index_version: v0.2-text-structure",
        "- scope: user-provided PPTX paths only",
        "- privacy_notice: This index contains absolute paths and extracted slide text.",
        "- storage_rule: Source PPTX files remain in place and are never overwritten.",
        f"- deck_count: {len(master)}",
        "",
        "## Decks",
        "",
    ]
    for path, index_path, deck_id, stat, slide_count, tags, titles in master:
        lines.extend(
            [
                f"### {path.stem}",
                f"- deck_id: `{deck_id}`",
                f"- source_path: `{path}`",
                f"- index_path: `{index_path}`",
                f"- slide_count: {slide_count}",
                f"- file_size_mb: {stat.st_size / 1024 / 1024:.1f}",
                f"- modified_time: {datetime.fromtimestamp(stat.st_mtime).isoformat(timespec='seconds')}",
                f"- deck_tags: {', '.join(tags)}",
                "- top_slide_titles:",
            ]
        )
        lines.extend(f"  - {title}" for title in titles)
        lines.append("")
    (output_dir / "MASTER_INDEX.md").write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", nargs="+", type=Path, help="User-provided PPTX paths.")
    parser.add_argument("--output-dir", type=Path, default=Path("PPT_Index"), help="Index output directory.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    paths = [path.expanduser().resolve() for path in args.pptx]
    build_index(paths, args.output_dir.expanduser().resolve())
    print(args.output_dir.expanduser().resolve())


if __name__ == "__main__":
    main()
