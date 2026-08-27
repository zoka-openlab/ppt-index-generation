from __future__ import annotations

import importlib.util
from pathlib import Path
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches


REPO_ROOT = Path(__file__).resolve().parents[1]
SCRIPT_DIR = REPO_ROOT / "skill" / "ppt-index-generation" / "scripts"


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


build_index_module = load_module("build_ppt_index", SCRIPT_DIR / "build_ppt_index.py")
extract_module = load_module("extract_candidate_slides", SCRIPT_DIR / "extract_candidate_slides.py")


def make_fictional_deck(path: Path, prefix: str = "Aurora") -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    titles = [
        f"{prefix} Studio Overview",
        "Product Portfolio",
        "Customer Outcomes",
    ]
    bodies = [
        "A fictional organization created only for automated testing.",
        "Planning, collaboration, and reporting capabilities.",
        "Faster reviews and clearer handoffs across sample teams.",
    ]
    for index, (title, body) in enumerate(zip(titles, bodies)):
        layout = prs.slide_layouts[0] if index == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
    table_slide = prs.slides[1]
    table = table_slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(1)).table
    table.cell(0, 0).text = "Capability"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Search"
    table.cell(1, 1).text = "Available"
    prs.save(path)
    return path


class SkillScriptTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)

    def tearDown(self):
        self.temp_dir.cleanup()

    def test_build_index_uses_synthetic_content_and_preserves_sources(self):
        source = make_fictional_deck(self.root / "source" / "sample.pptx")
        source_state = extract_module.file_state(source)
        output_dir = self.root / "private-index"

        build_index_module.build_index([source.resolve()], output_dir)

        self.assertEqual(extract_module.file_state(source), source_state)
        master = (output_dir / "MASTER_INDEX.md").read_text(encoding="utf-8")
        deck_indexes = list((output_dir / "indexes").glob("sample--*.md"))
        self.assertEqual(len(deck_indexes), 1)
        deck_index = deck_indexes[0].read_text(encoding="utf-8")
        self.assertIn("deck_count: 1", master)
        self.assertIn("slide_count: 3", master)
        self.assertIn("Aurora Studio Overview", deck_index)
        self.assertIn("Customer Outcomes", deck_index)
        self.assertIn("table", deck_index)

    def test_same_stem_decks_get_distinct_index_files(self):
        first = make_fictional_deck(self.root / "one" / "sample.pptx", "Aurora")
        second = make_fictional_deck(self.root / "two" / "sample.pptx", "Harbor")
        output_dir = self.root / "private-index"

        build_index_module.build_index([first.resolve(), second.resolve()], output_dir)

        deck_indexes = list((output_dir / "indexes").glob("sample--*.md"))
        self.assertEqual(len(deck_indexes), 2)

    def test_extract_preserves_requested_order_and_source(self):
        source = make_fictional_deck(self.root / "source.pptx")
        source_state = extract_module.file_state(source)
        output = self.root / "output" / "candidates.pptx"
        notes = self.root / "notes" / "candidates.md"

        extract_module.extract(source, [3, 1], output, notes, "Candidate pack")

        self.assertEqual(extract_module.file_state(source), source_state)
        result = Presentation(output)
        self.assertEqual(len(result.slides), 2)
        self.assertEqual(result.slides[0].shapes.title.text, "Customer Outcomes")
        self.assertEqual(result.slides[1].shapes.title.text, "Aurora Studio Overview")
        note_text = notes.read_text(encoding="utf-8")
        self.assertIn("source slide 3 - Customer Outcomes", note_text)
        self.assertIn("source slide 1 - Aurora Studio Overview", note_text)

    def test_extract_rejects_source_overwrite(self):
        source = make_fictional_deck(self.root / "source.pptx")

        with self.assertRaisesRegex(ValueError, "must not be the source"):
            extract_module.extract(source, [1], source, None, "Unsafe")

    def test_parse_slide_list_rejects_invalid_input(self):
        for value in ["", "0", "3-1", "a", "1,a"]:
            with self.subTest(value=value):
                with self.assertRaises(ValueError):
                    extract_module.parse_slide_list(value)

    def test_parse_slide_list_expands_ranges_and_deduplicates(self):
        self.assertEqual(extract_module.parse_slide_list("1,3-5,3"), [1, 3, 4, 5])


if __name__ == "__main__":
    unittest.main()
